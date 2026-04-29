from pptx import Presentation
import re

prs = Presentation('submission/SkyLoyal-Submission-template.pptx')

chrome_phrase = 'H9CEAI Customer Engagement & AI - MSCAIBUS1 - Pramithi Raroth Karimpanakkal'

total_all = 0
total_no_pl = 0
total_clean = 0

print(f"{'#':>2}  {'Title':<46}  {'all':>4}  {'-pl':>4}  {'clean':>5}")
print('-' * 72)

for i, slide in enumerate(prs.slides, 1):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    parts.append(run.text)
    text = ' '.join(parts)
    words_all = len(re.findall(r"\b\w+(?:'\w+)?\b", text))

    text_no_pl = re.sub(r'\[[^\]]*\]', '', text)
    words_no_pl = len(re.findall(r"\b\w+(?:'\w+)?\b", text_no_pl))

    text_clean = text_no_pl.replace(chrome_phrase, '')
    text_clean = re.sub(r'\b\d+\s*/\s*15\b', '', text_clean)
    words_clean = len(re.findall(r"\b\w+(?:'\w+)?\b", text_clean))

    title = ''
    for ln in [t.strip() for t in parts if t.strip()]:
        low = ln.lower()
        if 'h9ceai' in low or 'mscaibus' in low or '/ 15' in ln or ln.replace(' ', '').isdigit():
            continue
        title = ln[:46]
        break

    total_all += words_all
    total_no_pl += words_no_pl
    total_clean += words_clean
    print(f"{i:>2}  {title:<46}  {words_all:>4}  {words_no_pl:>4}  {words_clean:>5}")

print('-' * 72)
print(f"{'TOTAL':>50}  {total_all:>4}  {total_no_pl:>4}  {total_clean:>5}")
print()
print('all   = every word in the slide')
print('-pl   = minus [bracketed placeholders]')
print('clean = minus chrome/footer/page numbers (closest to Turnitin prose count)')
